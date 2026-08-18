from __future__ import annotations

import io
import json
import zipfile

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pypdf import PdfWriter

from app.services import knowledge_import as knowledge_import_module
from app.services.knowledge_import import (
    _OCR_RESULT_PREFIX,
    MAX_OCR_RENDER_PIXELS,
    KnowledgeImportError,
    _ocr_pdf_page_isolated,
    _pdf_ocr_render_scale,
    parse_payload,
    safe_file_name,
    validate_upload,
)


class _CompletedOcrProcess:
    returncode = 0

    def communicate(self, *, timeout: int):
        assert timeout > 0
        return (
            _OCR_RESULT_PREFIX
            + json.dumps({"text": "析境科技"}, ensure_ascii=False).encode("utf-8"),
            b"",
        )


def _docx_bytes(text: str) -> bytes:
    document = Document()
    document.add_paragraph(text)
    buffer = io.BytesIO()
    document.save(buffer)
    return buffer.getvalue()


def _pptx_bytes(text: str) -> bytes:
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[1]).shapes.title.text = text
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def _xlsx_bytes() -> bytes:
    workbook = Workbook()
    workbook.active.title = "产品"
    workbook.active.append(["名称", "说明"])
    workbook.active.append(["数智名片", "企业展示"])
    buffer = io.BytesIO()
    workbook.save(buffer)
    return buffer.getvalue()


def test_csv_import_produces_drafts_and_rejects_formula_injection() -> None:
    payload = "title,raw_text,visibility\n产品说明,正文内容,internal\n".encode()
    assert validate_upload("docs.csv", "text/csv", payload) == "csv"
    assert parse_payload("csv", "docs.csv", payload)[0].visibility == "internal"

    with pytest.raises(KnowledgeImportError, match="IMPORT_DANGEROUS_VALUE"):
        parse_payload("csv", "docs.csv", b"title,raw_text\n=cmd,content\n")


def test_csv_import_accepts_tabular_columns_and_rejects_oversized_cells() -> None:
    draft = parse_payload("csv", "docs.csv", "title,raw_text,category\nA,content,产品\n".encode())[
        0
    ]
    assert "category: 产品" in draft.raw_text

    oversized = b"raw_text\n" + (b"a" * 100_001) + b"\n"
    with pytest.raises(KnowledgeImportError, match="IMPORT_CSV_CELL_TOO_LARGE"):
        parse_payload("csv", "docs.csv", oversized)


def test_docx_is_parsed_but_macro_and_archive_paths_are_rejected() -> None:
    payload = _docx_bytes("企业安全知识")
    assert (
        validate_upload(
            "guide.docx",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            payload,
        )
        == "docx"
    )
    assert parse_payload("docx", "guide.docx", payload)[0].raw_text == "企业安全知识"

    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w") as archive:
        archive.writestr("word/document.xml", "<document/>")
        archive.writestr("../vbaProject.bin", b"macro")
    with pytest.raises(KnowledgeImportError, match="IMPORT_ARCHIVE_PATH"):
        validate_upload(
            "bad.docx",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            buffer.getvalue(),
        )


def test_encrypted_pdf_and_mime_magic_mismatches_are_rejected() -> None:
    writer = PdfWriter()
    writer.add_blank_page(width=100, height=100)
    writer.encrypt("secret")
    buffer = io.BytesIO()
    writer.write(buffer)
    with pytest.raises(KnowledgeImportError, match="IMPORT_ENCRYPTED_PDF"):
        parse_payload("pdf", "secret.pdf", buffer.getvalue())


def test_pdf_layout_controls_are_normalized_but_other_controls_stay_rejected(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    class Page:
        def __init__(self, text: str) -> None:
            self._text = text

        def extract_text(self) -> str:
            return self._text

    class Reader:
        is_encrypted = False

        def __init__(self, text: str) -> None:
            self.pages = [Page(text)]

    safe_text = ("项目背景与解决方案" * 12) + "\x0b分页\x0c成果"
    monkeypatch.setattr(
        knowledge_import_module,
        "PdfReader",
        lambda *_args, **_kwargs: Reader(safe_text),
    )
    draft = knowledge_import_module._parse_pdf("case.pdf", b"pdf")
    assert "\x0b" not in draft.raw_text
    assert "\x0c" not in draft.raw_text

    dangerous_text = ("项目背景与解决方案" * 12) + "\x01危险"
    monkeypatch.setattr(
        knowledge_import_module,
        "PdfReader",
        lambda *_args, **_kwargs: Reader(dangerous_text),
    )
    with pytest.raises(KnowledgeImportError, match="IMPORT_DANGEROUS_VALUE"):
        knowledge_import_module._parse_pdf("case.pdf", b"pdf")

    with pytest.raises(KnowledgeImportError, match="IMPORT_MIME_MISMATCH"):
        validate_upload("file.pdf", "text/plain", b"%PDF-1.7")
    with pytest.raises(KnowledgeImportError, match="IMPORT_MAGIC_MISMATCH"):
        validate_upload("file.pdf", "application/pdf", b"not-a-pdf")


def test_pdf_ocr_rendering_caps_exported_slide_pixel_budget() -> None:
    assert _pdf_ocr_render_scale(595, 842) == 2.0

    width = 3_840
    height = 2_160
    scale = _pdf_ocr_render_scale(width, height)
    assert width * scale * height * scale == pytest.approx(MAX_OCR_RENDER_PIXELS)

    with pytest.raises(KnowledgeImportError, match="IMPORT_PDF_INVALID"):
        _pdf_ocr_render_scale(0, height)


def test_pdf_ocr_page_runs_in_disposable_process(monkeypatch: pytest.MonkeyPatch) -> None:
    captured: dict[str, object] = {}

    def fake_popen(command, **kwargs):
        captured["command"] = command
        captured["kwargs"] = kwargs
        return _CompletedOcrProcess()

    monkeypatch.setattr("app.services.knowledge_import.subprocess.Popen", fake_popen)
    assert _ocr_pdf_page_isolated("source.pdf", page_number=2) == "析境科技"
    assert captured["command"][-4:] == [
        "--pdf-path",
        "source.pdf",
        "--page-number",
        "2",
    ]


def test_office_and_html_formats_extract_text_without_network_access() -> None:
    pptx = _pptx_bytes("企业介绍")
    assert (
        validate_upload(
            "intro.pptx",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            pptx,
        )
        == "pptx"
    )
    assert "企业介绍" in parse_payload("pptx", "intro.pptx", pptx)[0].raw_text

    xlsx = _xlsx_bytes()
    assert (
        validate_upload(
            "catalog.xlsx",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            xlsx,
        )
        == "xlsx"
    )
    assert "数智名片" in parse_payload("xlsx", "catalog.xlsx", xlsx)[0].raw_text

    html = (
        b"<h1>\xe4\xbc\x81\xe4\xb8\x9a</h1><script>ignore()</script><p>\xe4\xba\xa7\xe5\x93\x81</p>"
    )
    assert "企业" in parse_payload("html", "page.html", html)[0].raw_text
    assert "产品" in parse_payload("html", "page.html", html)[0].raw_text
    assert "ignore" not in parse_payload("html", "page.html", html)[0].raw_text


def test_extracted_document_controls_are_normalized_instead_of_rejected() -> None:
    draft = parse_payload("txt", "report.txt", b"first page\x0csecond page\x00end")[0]

    assert draft.raw_text == "first page\nsecond page\nend"
    assert "\x0c" not in draft.raw_text
    assert "\x00" not in draft.raw_text


def test_upload_validation_has_no_application_file_size_limit() -> None:
    payload = b"x" * (10 * 1024 * 1024 + 1)

    assert validate_upload("large.txt", "text/plain", payload) == "txt"


@pytest.mark.parametrize("name", ["../file.csv", "folder/file.csv", "folder\\file.csv", ""])
def test_unsafe_file_names_are_rejected(name: str) -> None:
    with pytest.raises(KnowledgeImportError, match="IMPORT_UNSAFE_FILENAME"):
        safe_file_name(name)
